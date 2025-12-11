"""
Demo script to test the bullet_style feature in python-pptx.

This script demonstrates how to use the new bullet_style property
on paragraphs to create bulleted and numbered lists.

Usage:
    cd python-pptx
    python src/demo_bullet_style.py
    
    Open the generated file in PowerPoint to verify bullet styles.
"""

import os
import sys
sys.path.insert(0, 'src')

from pptx import Presentation
from pptx.util import Inches, Pt

# Output directory for generated files
OUTPUT_DIR = 'src/demo_test_results'


def ensure_output_dir():
    """Create output directory if it doesn't exist."""
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
        print(f"Created output directory: {OUTPUT_DIR}")
    return OUTPUT_DIR


# Ensure output directory exists
ensure_output_dir()

# Create a new presentation
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Add a text box
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(5)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# First paragraph - no bullet (default)
p1 = tf.paragraphs[0]
p1.text = "This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default)"
p1.font.size = Pt(7)
print(f"Paragraph 1 - No bullet: bullet_style = {p1.bullet_style}")

# Second paragraph - bullet style (as requested in the example)
p2 = tf.add_paragraph()
p2.text = "$XXk annual software spend, split by XX per segment This is a paragraph without any bullet style (default) This is a paragraph without any bu This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default)llet style (default) This is a paragraph without any bullet style (default)"
p2.level = 0
p2.bullet_style = "bullet"
p2.space_before = Pt(3)  # Add spacing before bullet items
p2.font.size = Pt(7)
print(f"Paragraph 2 - Bullet: bullet_style = {p2.bullet_style}")

# Third paragraph - another bullet
p3 = tf.add_paragraph()
p3.text = "Cost breakdown by department This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default)"
p3.level = 0
p3.bullet_style = "bullet"
p3.space_before = Pt(3)  # Add spacing before bullet items
p3.font.size = Pt(7)
print(f"Paragraph 3 - Bullet: bullet_style = {p3.bullet_style}")

# Third paragraph - another bullet
p3i = tf.add_paragraph()
p3i.text = "INDENTED cost breakdown by department This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default)"
p3i.level = 2
p3i.bullet_style = "bullet"
p3i.space_before = Pt(3)  # Add spacing before bullet items
p3i.font.size = Pt(7)
print(f"Paragraph 3 Indented - Bullet: bullet_style = {p3i.bullet_style}")

# Fourth paragraph - numbered style
p4 = tf.add_paragraph()
p4.text = "First action item This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default) This is a paragraph without any bullet style (default)"
p4.level = 0
p4.bullet_style = "number"
p4.space_before = Pt(3)  # Add spacing before numbered items
p4.font.size = Pt(7)
print(f"Paragraph 4 - Number: bullet_style = {p4.bullet_style}")

# Fifth paragraph - also numbered (continues numbering)
p5 = tf.add_paragraph()
p5.text = "Second action item"
p5.level = 0
p5.bullet_style = "number"
p5.space_before = Pt(3)  # Add spacing before numbered items
p5.font.size = Pt(7)
print(f"Paragraph 5 - Number: bullet_style = {p5.bullet_style}")

# Sixth paragraph - Level 1 nested bullet (demonstrates level-based indentation)
p6 = tf.add_paragraph()
p6.text = "Level 1 sub-item - notice increased indentation"
p6.level = 1
p6.bullet_style = "bullet"
p6.space_before = Pt(3)
p6.font.size = Pt(7)
print(f"Paragraph 6 - Level 1 bullet: bullet_style = {p6.bullet_style}, level = {p6.level}")

# Seventh paragraph - Level 2 nested bullet
p7 = tf.add_paragraph()
p7.text = "Level 2 sub-item - even more indentation"
p7.level = 2
p7.bullet_style = "bullet"
p7.space_before = Pt(3)
p7.font.size = Pt(7)
print(f"Paragraph 7 - Level 2 bullet: bullet_style = {p7.bullet_style}, level = {p7.level}")

# Eighth paragraph - Back to Level 0
p8 = tf.add_paragraph()
p8.text = ""
p8.level = 0
p8.bullet_style = "bullet"
p8.space_before = Pt(3)
p8.font.size = Pt(7)
print(f"Paragraph 8 - Level 0 bullet: bullet_style = {p8.bullet_style}, level = {p8.level}")

# Ninth paragraph - no bullet (default)
p9 = tf.add_paragraph()
p9.text = "This is a paragraph without any bullet style (default)"
p9.font.size = Pt(7)
print(f"Paragraph 9 - No bullet: bullet_style = {p9.bullet_style}")


# Save the presentation
output_path = os.path.join(OUTPUT_DIR, "bullet_style_demo.pptx")
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print("Open the file in PowerPoint to verify the bullet styles are correct.")

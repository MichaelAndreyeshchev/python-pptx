"""Inspect how bullet_style reads existing bullets from PresentationWBullets.pptx"""

import sys
sys.path.insert(0, 'src')
from pptx import Presentation
from lxml import etree

# Open the existing presentation with bullets
prs = Presentation('src/PresentationWBullets.pptx')

print('=== Inspecting PresentationWBullets.pptx ===')
print(f'Total slides: {len(prs.slides)}')
print()

for slide_idx, slide in enumerate(prs.slides):
    print(f'=== Slide {slide_idx + 1} ===')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        tf = shape.text_frame
        print(f'  Shape: "{shape.name}"')
        
        for p_idx, p in enumerate(tf.paragraphs):
            # KEY: This shows how bullet_style reads existing bullets
            bullet_style = p.bullet_style
            level = p.level
            text_preview = p.text[:50] + '...' if len(p.text) > 50 else p.text
            
            # Also show the underlying XML elements
            pPr = p._p.pPr
            has_buChar = pPr.buChar if pPr is not None else None
            has_buAutoNum = pPr.buAutoNum if pPr is not None else None
            has_buNone = pPr.buNone if pPr is not None else None
            
            print(f'    Para {p_idx}: level={level}, bullet_style="{bullet_style}", text="{text_preview}"')
            if has_buChar is not None:
                print(f'             XML: buChar exists (char="{has_buChar.char}")')
            if has_buAutoNum is not None:
                print(f'             XML: buAutoNum exists (type="{has_buAutoNum.type}")')
            if has_buNone is not None:
                print(f'             XML: buNone exists')
        print()

print("\n" + "="*80)
print("=== Raw XML for paragraphs with bullets ===")
print("="*80 + "\n")

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        tf = shape.text_frame
        for p_idx, p in enumerate(tf.paragraphs):
            if p.bullet_style is not None:
                print(f'Shape "{shape.name}", Paragraph {p_idx}:')
                # Get the XML for this paragraph
                xml_str = etree.tostring(p._p, pretty_print=True, encoding='unicode')
                print(xml_str)
                print()


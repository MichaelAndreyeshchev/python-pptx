"""Inspect the structure of image_exercise.pptx"""

from pptx import Presentation

# Open the existing presentation
prs = Presentation('image_exercise.pptx')

print(f'Total slides: {len(prs.slides)}')
print()

for slide_idx, slide in enumerate(prs.slides):
    print(f'=== Slide {slide_idx + 1} ===')
    for shape_idx, shape in enumerate(slide.shapes):
        print(f'  Shape {shape_idx}: {shape.shape_type}, name="{shape.name}"')
        if shape.has_text_frame:
            print(f'    Has text frame with {len(shape.text_frame.paragraphs)} paragraph(s)')
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                text_preview = p.text[:50] + '...' if len(p.text) > 50 else p.text
                print(f'      P{p_idx}: "{text_preview}"')
        if hasattr(shape, 'left'):
            print(f'    Position: left={shape.left}, top={shape.top}')
    print()


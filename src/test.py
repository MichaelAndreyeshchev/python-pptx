from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
tf = txBox.text_frame

# Create a bullet point
p1 = tf.add_paragraph()
p1.text = "$XXk annual software spend, split by XX per segment"
p1.level = 0
p1.bullet_style = "bullet"

# Create a numbered item
p2 = tf.add_paragraph()
p2.text = "First action item"
p2.level = 0
p2.bullet_style = "number"

# Create a numbered item
p3 = tf.add_paragraph()
p3.text = "Second action item"
p3.level = 0
p3.bullet_style = "number"

# Remove bullet formatting
p3.bullet_style = None

prs.save("output.pptx")
"""
Test script demonstrating bullet_style on an existing PowerPoint presentation.

This script opens image_exercise.pptx and adds:
- Bullet points
- Numbered lists  
- Regular text entries

to various text boxes in the presentation.
"""

import os
import sys
sys.path.insert(0, 'src')

from pptx import Presentation
from pptx.util import Pt

# Output directory for generated files
OUTPUT_DIR = 'src/demo_test_results'


def ensure_output_dir():
    """Create output directory if it doesn't exist."""
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
        print(f"Created output directory: {OUTPUT_DIR}")
    return OUTPUT_DIR


def main():
    # Ensure output directory exists
    ensure_output_dir()
    
    # Open the existing presentation
    print("Opening image_exercise.pptx...")
    prs = Presentation('src/image_exercise.pptx')
    
    slide = prs.slides[0]
    
    # Track which shapes we modified
    modifications = []
    
    # Find text boxes and add different content types
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        tf = shape.text_frame

        print(shape.name)
        
        if "Textfeld 183" in shape.name: 
            tf.clear()  # Clear existing content
            
            p1 = tf.paragraphs[0]
            p1.text = "Key Features"
            p1.font.size = Pt(7)
            p1.font.bold = True
            
            # Bullet points
            p2 = tf.add_paragraph()
            p2.text = "Cloud-native architecture " + shape.name
            p2.bullet_style = "number"
            p2.font.size = Pt(7)
            p2.space_before = Pt(3)
            
            p3 = tf.add_paragraph()
            p3.text = "Real-time analytics"
            p3.level = 1
            p3.bullet_style = "number"
            p3.font.size = Pt(7)
            p3.space_before = Pt(3)
            
            p4 = tf.add_paragraph()
            p4.text = "Enterprise security"
            p4.level = 2
            p4.bullet_style = "number"
            p4.font.size = Pt(7)
            p4.space_before = Pt(3)

            p5 = tf.add_paragraph()
            p5.text = "Speed"
            p5.level = 2
            p5.bullet_style = "number"
            p5.font.size = Pt(7)
            p5.space_before = Pt(3)  

        if "Textfeld 170" in shape.name: 
            tf.clear()  # Clear existing content
            
            p1 = tf.paragraphs[0]
            p1.text = "Key Features"
            p1.font.size = Pt(7)
            p1.font.bold = True
            
            # Bullet points
            p2 = tf.add_paragraph()
            p2.text = "Cloud-native architecture " + shape.name
            p2.bullet_style = "bullet"
            p2.font.size = Pt(7)
            p2.space_before = Pt(3)

            p2 = tf.add_paragraph()
            p2.text = "test"
            p2.bullet_style = "bullet"
            p2.level = 1
            p2.font.size = Pt(7)
            p2.space_before = Pt(3)
            
            p3 = tf.add_paragraph()
            p3.text = "Real-time analytics"
            p3.bullet_style = "number"
            p3.font.size = Pt(7)
            p3.space_before = Pt(5)
            
            p4 = tf.add_paragraph()
            p4.text = "Enterprise security"
            p4.bullet_style = None
            p4.font.size = Pt(7)
            p4.space_before = Pt(3) 

    # Save the modified presentation
    output_path = os.path.join(OUTPUT_DIR, "image_exercise_with_bullets.pptx")
    prs.save(output_path)
    
    # Print summary
    print(f"\n=== Modifications Made ===")
    for mod in modifications:
        print(f"  [OK] {mod}")
    
    print(f"\nSaved modified presentation to: {output_path}")
    print("\nOpen the file in PowerPoint to verify:")
    print("  - Bullet points (•) with proper indentation")
    print("  - Numbered lists (1., 2., 3.) with proper indentation")
    print("  - Regular text paragraphs")
    print("  - Proper spacing between items")


def test_bullet_style_reading():
    """Test that bullet_style correctly reads existing bullets from the modified file."""
    print("\n=== Testing bullet_style reading ===")
    
    output_path = os.path.join(OUTPUT_DIR, "image_exercise_with_bullets.pptx")
    try:
        prs = Presentation(output_path)
    except FileNotFoundError:
        print(f"Run main() first to create the modified file at {output_path}")
        return
    
    slide = prs.slides[0]
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        tf = shape.text_frame
        has_bullets = False
        
        for p in tf.paragraphs:
            if p.bullet_style is not None:
                has_bullets = True
                break
        
        if has_bullets:
            print(f"\nShape: {shape.name}")
            for i, p in enumerate(tf.paragraphs):
                style = p.bullet_style if p.bullet_style else "none"
                text_preview = p.text[:40] + "..." if len(p.text) > 40 else p.text
                print(f"  P{i}: style={style:8s} text=\"{text_preview}\"")


if __name__ == "__main__":
    main()
    test_bullet_style_reading()


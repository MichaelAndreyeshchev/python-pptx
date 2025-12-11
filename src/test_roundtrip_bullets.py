"""
Test script: Read existing bullets from PPTX, then save and verify.

This demonstrates the full round-trip:
1. Open existing PPTX with bullets
2. Read and display bullet_style values
3. Save to a new file
4. Re-open and verify bullets are preserved
"""

import sys
import os
sys.path.insert(0, 'src')

from pptx import Presentation
from pptx.util import Pt, Inches

# Output directory for generated files
OUTPUT_DIR = 'src/demo_test_results'


def ensure_output_dir():
    """Create output directory if it doesn't exist."""
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
        print(f"Created output directory: {OUTPUT_DIR}")
    return OUTPUT_DIR


def print_bullets(prs, label=""):
    """Print all bullet information from a presentation."""
    print(f"\n{'='*60}")
    print(f"  {label}")
    print(f"{'='*60}")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nSlide {slide_idx + 1}:")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            has_bullets = any(p.bullet_style is not None for p in tf.paragraphs)
            
            if has_bullets or any(p.text.strip() for p in tf.paragraphs):
                print(f'  Shape: "{shape.name}"')
                for p_idx, p in enumerate(tf.paragraphs):
                    bullet = p.bullet_style if p.bullet_style else "none"
                    level = p.level
                    text = p.text[:40] + "..." if len(p.text) > 40 else p.text
                    print(f'    [{p_idx}] bullet_style="{bullet}", level={level}, text="{text}"')


def test_read_existing_bullets():
    """Test 1: Read bullet styles from existing PPTX."""
    print("\n" + "="*60)
    print("  TEST 1: Reading bullets from existing PPTX")
    print("="*60)
    
    # Open the existing presentation
    input_file = 'src/PresentationWBullets.pptx'
    print(f"\nOpening: {input_file}")
    
    prs = Presentation(input_file)
    print_bullets(prs, "Original File Contents")
    
    return prs


def test_save_and_reload(prs):
    """Test 2: Save presentation and reload to verify bullets persist."""
    print("\n" + "="*60)
    print("  TEST 2: Save and reload (round-trip test)")
    print("="*60)
    
    # Save to a new file
    output_file = os.path.join(OUTPUT_DIR, 'PresentationWBullets_roundtrip.pptx')
    print(f"\nSaving to: {output_file}")
    prs.save(output_file)
    print("  [OK] Saved successfully")
    
    # Reload and verify
    print(f"\nReloading: {output_file}")
    prs_reloaded = Presentation(output_file)
    print_bullets(prs_reloaded, "Reloaded File Contents")
    
    # Verify bullets match
    print("\n" + "-"*40)
    print("  VERIFICATION: Comparing original vs reloaded")
    print("-"*40)
    
    original_prs = Presentation('src/PresentationWBullets.pptx')
    
    all_match = True
    for slide_idx, (orig_slide, reload_slide) in enumerate(zip(original_prs.slides, prs_reloaded.slides)):
        for orig_shape, reload_shape in zip(orig_slide.shapes, reload_slide.shapes):
            if not orig_shape.has_text_frame:
                continue
            
            for p_idx, (orig_p, reload_p) in enumerate(zip(orig_shape.text_frame.paragraphs, 
                                                           reload_shape.text_frame.paragraphs)):
                if orig_p.bullet_style != reload_p.bullet_style:
                    print(f"  [MISMATCH] Slide {slide_idx+1}, Shape '{orig_shape.name}', Para {p_idx}")
                    print(f"             Original: {orig_p.bullet_style}, Reloaded: {reload_p.bullet_style}")
                    all_match = False
    
    if all_match:
        print("  [OK] All bullet styles preserved correctly!")
    
    return prs_reloaded


def test_modify_bullets():
    """Test 3: Modify bullets and save."""
    print("\n" + "="*60)
    print("  TEST 3: Modify bullets and create new content")
    print("="*60)
    
    # Start fresh from original
    prs = Presentation('src/PresentationWBullets.pptx')
    slide = prs.slides[0]
    
    # Find the shape with bullets and modify
    for shape in slide.shapes:
        if shape.name == "Subtitle 2" and shape.has_text_frame:
            tf = shape.text_frame
            
            print(f"\nModifying shape: '{shape.name}'")
            print("  Before:")
            for p_idx, p in enumerate(tf.paragraphs):
                print(f"    [{p_idx}] bullet_style='{p.bullet_style}', text='{p.text}'")
            
            # Add a new paragraph with numbered bullet
            new_p = tf.add_paragraph()
            new_p.text = "New numbered item"
            new_p.bullet_style = "number"
            new_p.font.size = Pt(18)
            
            # Add another bullet point
            new_p2 = tf.add_paragraph()
            new_p2.text = "Another bullet point"
            new_p2.bullet_style = "bullet"
            new_p2.font.size = Pt(18)
            
            # Add indented bullet (level 1)
            new_p3 = tf.add_paragraph()
            new_p3.text = "Indented sub-item"
            new_p3.level = 1
            new_p3.bullet_style = "bullet"
            new_p3.font.size = Pt(16)

            # Add indented bullet (level 2)
            new_p4 = tf.add_paragraph()
            new_p4.text = "Another indented sub-item"
            new_p4.level = 2
            new_p4.bullet_style = "bullet"
            new_p4.font.size = Pt(16)
            
            # Add plain text (no bullet)
            new_p6 = tf.add_paragraph()
            new_p6.text = "Plain text - no bullet"
            new_p6.bullet_style = None  # Explicitly no bullet
            new_p6.font.size = Pt(18)
            
            print("\n  After:")
            for p_idx, p in enumerate(tf.paragraphs):
                print(f"    [{p_idx}] bullet_style='{p.bullet_style}', level={p.level}, text='{p.text}'")
    
    # Save modified version
    output_file = os.path.join(OUTPUT_DIR, 'PresentationWBullets_modified.pptx')
    print(f"\nSaving modified version to: {output_file}")
    prs.save(output_file)
    print("  [OK] Saved successfully")
    
    # Verify by reloading
    print(f"\nVerifying by reloading: {output_file}")
    prs_verify = Presentation(output_file)
    print_bullets(prs_verify, "Modified & Reloaded Contents")
    
    return output_file


def test_create_from_scratch():
    """Test 4: Create presentation with bullets from scratch."""
    print("\n" + "="*60)
    print("  TEST 4: Create new presentation with bullets")
    print("="*60)
    
    # Create new presentation
    prs = Presentation()
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # Title (no bullet)
    p0 = tf.paragraphs[0]
    p0.text = "My Bullet List Demo"
    p0.font.size = Pt(24)
    p0.font.bold = True
    
    # Bullet items
    items = [
        ("First bullet point", "bullet", 0),
        ("Second bullet point", "bullet", 0),
        ("Sub-item under second", "bullet", 1),
        ("Another sub-item", "bullet", 2),
        ("Another sub-item", "bullet", 1),
        ("Back to main level", "bullet", 0),
        ("First numbered item", "number", 0),
        ("Second numbered item", "number", 1),
        ("Third numbered item", "number", 2),
        ("Plain text paragraph", None, 0),
    ]
    
    for text, style, level in items:
        p = tf.add_paragraph()
        p.text = text
        p.bullet_style = style
        p.level = level
        p.font.size = Pt(18) if level == 0 else Pt(16)
        p.space_before = Pt(6)
    
    print("\nCreated paragraphs:")
    for p_idx, p in enumerate(tf.paragraphs):
        print(f"  [{p_idx}] bullet_style='{p.bullet_style}', level={p.level}, text='{p.text}'")
    
    # Save
    output_file = os.path.join(OUTPUT_DIR, 'PresentationWBullets_created.pptx')
    print(f"\nSaving to: {output_file}")
    prs.save(output_file)
    print("  [OK] Saved successfully")
    
    # Verify
    print(f"\nVerifying by reloading: {output_file}")
    prs_verify = Presentation(output_file)
    print_bullets(prs_verify, "Created & Reloaded Contents")
    
    return output_file


def main():
    """Run all tests."""
    print("\n" + "#"*60)
    print("#  BULLET STYLE ROUND-TRIP TEST SUITE")
    print("#"*60)
    
    # Ensure output directory exists
    ensure_output_dir()
    
    # Test 1: Read existing
    prs = test_read_existing_bullets()
    
    # Test 2: Save and reload
    test_save_and_reload(prs)
    
    # Test 3: Modify and save
    test_modify_bullets()
    
    # Test 4: Create from scratch
    test_create_from_scratch()
    
    print("\n" + "#"*60)
    print("#  ALL TESTS COMPLETE")
    print("#"*60)
    print(f"\nGenerated files in '{OUTPUT_DIR}/':")
    print("  - PresentationWBullets_roundtrip.pptx  (exact copy)")
    print("  - PresentationWBullets_modified.pptx   (with added bullets)")
    print("  - PresentationWBullets_created.pptx    (created from scratch)")
    print("\nOpen these in PowerPoint to verify the bullets display correctly!")


if __name__ == "__main__":
    main()


"""Unit tests for the bullet_style property on _Paragraph."""

import pytest

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.text.text import _Paragraph  # type: ignore[reportPrivateUsage]


class DescribeParagraphBulletStyle:
    """Unit tests for _Paragraph.bullet_style property."""

    def it_returns_None_when_no_bullet_element_is_present(self):
        p_xml = '<a:p %s><a:pPr/></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        assert paragraph.bullet_style is None

    def it_returns_None_when_pPr_is_not_present(self):
        p_xml = '<a:p %s/>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        assert paragraph.bullet_style is None

    def it_returns_bullet_when_buChar_element_is_present(self):
        p_xml = '<a:p %s><a:pPr><a:buChar char="•"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        assert paragraph.bullet_style == "bullet"

    def it_returns_number_when_buAutoNum_element_is_present(self):
        p_xml = '<a:p %s><a:pPr><a:buAutoNum type="arabicPeriod"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        assert paragraph.bullet_style == "number"

    def it_returns_None_when_buNone_element_is_present(self):
        p_xml = '<a:p %s><a:pPr><a:buNone/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        assert paragraph.bullet_style is None

    def it_can_set_bullet_style_to_bullet(self):
        p_xml = '<a:p %s/>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = "bullet"
        
        assert paragraph.bullet_style == "bullet"
        pPr = p.pPr
        assert pPr is not None
        assert pPr.buChar is not None
        assert pPr.buChar.char == "\u2022"

    def it_can_set_bullet_style_to_number(self):
        p_xml = '<a:p %s/>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = "number"
        
        assert paragraph.bullet_style == "number"
        pPr = p.pPr
        assert pPr is not None
        assert pPr.buAutoNum is not None
        assert pPr.buAutoNum.type == "arabicPeriod"

    def it_can_remove_bullet_style_by_setting_to_None(self):
        p_xml = '<a:p %s><a:pPr><a:buChar char="•"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = None
        
        assert paragraph.bullet_style is None
        pPr = p.pPr
        assert pPr.buChar is None

    def it_can_remove_bullet_style_by_setting_to_empty_string(self):
        p_xml = '<a:p %s><a:pPr><a:buAutoNum type="arabicPeriod"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = ""
        
        assert paragraph.bullet_style is None
        pPr = p.pPr
        assert pPr.buAutoNum is None

    def it_can_change_bullet_style_from_bullet_to_number(self):
        p_xml = '<a:p %s><a:pPr><a:buChar char="•"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = "number"
        
        assert paragraph.bullet_style == "number"
        pPr = p.pPr
        assert pPr.buChar is None
        assert pPr.buAutoNum is not None

    def it_can_change_bullet_style_from_number_to_bullet(self):
        p_xml = '<a:p %s><a:pPr><a:buAutoNum type="arabicPeriod"/></a:pPr></a:p>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        paragraph.bullet_style = "bullet"
        
        assert paragraph.bullet_style == "bullet"
        pPr = p.pPr
        assert pPr.buAutoNum is None
        assert pPr.buChar is not None

    def it_raises_ValueError_for_invalid_bullet_style(self):
        p_xml = '<a:p %s/>' % nsdecls('a')
        p = parse_xml(p_xml)
        paragraph = _Paragraph(p, None)
        
        with pytest.raises(ValueError) as exc_info:
            paragraph.bullet_style = "invalid"
        
        assert "bullet_style must be None, '', 'bullet', or 'number'" in str(exc_info.value)


class DescribeIntegrationBulletStyle:
    """Integration tests for bullet_style using real Presentation objects."""

    def it_works_with_real_presentation(self):
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        
        p1 = tf.paragraphs[0]
        p1.text = "No bullet"
        assert p1.bullet_style is None
        
        p2 = tf.add_paragraph()
        p2.text = "Bullet item"
        p2.bullet_style = "bullet"
        assert p2.bullet_style == "bullet"
        
        p3 = tf.add_paragraph()
        p3.text = "Numbered item"
        p3.bullet_style = "number"
        assert p3.bullet_style == "number"
